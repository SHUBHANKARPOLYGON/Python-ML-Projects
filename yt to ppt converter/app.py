from flask import Flask, render_template, request, send_file
from pytube import YouTube
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from pptx.util import Inches
import os
import google.generativeai as genai

# Configure Gemini
genai.configure(api_key="AIzaSyAjgZwy5SMT-nNjfutVpMu7N7WGCclFFrs")  # Replace with your actual key

app = Flask(__name__)

# === Utility Functions ===

def get_transcript(video_url):
    video_id = YouTube(video_url).video_id
    transcript = YouTubeTranscriptApi.get_transcript(video_id)
    full_text = " ".join([item['text'] for item in transcript])
    return full_text

def gemini_summary(text):
    model = genai.GenerativeModel(model_name="gemini-pro")
    prompt = f"Summarize the following YouTube transcript into 5 slides, each with 3-5 bullet points:\n\n{text}"
    response = model.generate_content(prompt)
    return response.text.split("\n\n")  # Split into slide-sized chunks

def create_ppt(slide_texts, output_path="output.pptx"):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "YouTube Video Summary"
    slide.placeholders[1].text = "Generated with Gemini AI"

    # Content slides
    for i, text in enumerate(slide_texts):
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = text.strip()

    prs.save(output_path)

# === Routes ===

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        url = request.form["video_url"]
        try:
            transcript = get_transcript(url)
            slides = gemini_summary(transcript)
            create_ppt(slides, "output.pptx")
            return send_file("output.pptx", as_attachment=True)
        except Exception as e:
            return f"Error: {str(e)}"
    return render_template("index.html")

if __name__ == "__main__":
    app.run(debug=True)
